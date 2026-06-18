import datetime
import math
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import tools.chart_generator as cg



def _get(item, key, default=""):
    if isinstance(item, dict):
        return item.get(key, default)
    return getattr(item, key, default)



def _fmt_number(value, digits=2):
    try:
        return f"{float(value):.{digits}f}"
    except (TypeError, ValueError):
        return "N/A" if value is None else str(value)



def _to_float(value):
    try:
        return float(value)
    except (TypeError, ValueError):
        return None



def _shorten(text, limit=110):
    text = str(text or "").strip()
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "…"


_URL_RE = re.compile(r"https?://\S+", re.IGNORECASE)
_PARENS_URL_RE = re.compile(r"\([^)]*https?://[^)]*\)", re.IGNORECASE)
_TIMELINE_SUMMARY_FALLBACK = ""
_CJK_RE = re.compile(r"[\u4e00-\u9fff]")
_TIMELINE_SUMMARY_PATCH_RE = re.compile(
    r"^(进一步看|更进一步|补充判断：围绕|补充判断|值得持续关注|后续仍需关注|"
    r"该事件具有重要意义|相关进展值得跟踪)[，,:：]?.*$",
    re.IGNORECASE,
)


def _clean_display_text(text):
    cleaned = str(text or "").strip()
    cleaned = _PARENS_URL_RE.sub("", cleaned)
    cleaned = _URL_RE.sub("", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned)
    return cleaned.strip(" -:：;；,，")


def _format_timeline_event_summary(item):
    summary = _clean_display_text(_get(item, "event_summary", ""))
    summary = re.sub(r"\[\s*(?:\.\.\.|…)+\s*\]", " ", summary).strip()
    if not summary or _TIMELINE_SUMMARY_PATCH_RE.match(summary):
        return _TIMELINE_SUMMARY_FALLBACK
    cjk_count = len(_CJK_RE.findall(summary))
    if cjk_count < 10 or cjk_count / max(len(summary), 1) < 0.35:
        return _TIMELINE_SUMMARY_FALLBACK
    return _shorten(summary, 220)


def _estimate_line_units(text, chars_per_line=34):
    weighted_length = 0.0
    for ch in str(text or ""):
        if "\u4e00" <= ch <= "\u9fff":
            weighted_length += 1.0
        elif ord(ch) < 128:
            weighted_length += 0.55
        else:
            weighted_length += 0.8
    return max(1, int(math.ceil(weighted_length / max(chars_per_line, 1))))


def _make_entry(text, font_size=13, bold=False, color=None, space_after=6, hyperlink=None):
    return {
        "text": str(text or "").strip(),
        "font_size": font_size,
        "bold": bold,
        "color": color,
        "space_after": space_after,
        "hyperlink": hyperlink,
    }


_SECTION_LINE_RE = re.compile(r"^(【[^】]+】)\s*(.*)$")


def _append_summary_entries(entries, summary_text, body_font=13, space_after=6):
    for raw_line in str(summary_text or "").split("\n"):
        line = raw_line.strip()
        if not line:
            continue
        match = _SECTION_LINE_RE.match(line)
        if match:
            section_title, section_body = match.groups()
            entries.append(
                _make_entry(
                    section_title,
                    font_size=body_font,
                    bold=True,
                    color=(0, 51, 102),
                    space_after=2,
                )
            )
            if section_body.strip():
                entries.append(
                    _make_entry(
                        section_body.strip(),
                        font_size=body_font,
                        bold=False,
                        color=(32, 32, 32),
                        space_after=space_after,
                    )
                )
            continue

        entries.append(
            _make_entry(
                line,
                font_size=body_font,
                bold=False,
                color=(32, 32, 32),
                space_after=space_after,
            )
        )


def _entry_units(entry, chars_per_line=34):
    text = str(entry.get("text", "") or "").strip()
    if not text:
        return 0
    units = _estimate_line_units(text, chars_per_line=chars_per_line)
    units += 1 if int(entry.get("font_size", 13) or 13) >= 14 else 0
    return units


def _total_units(entries, chars_per_line=34):
    return sum(_entry_units(entry, chars_per_line=chars_per_line) for entry in entries)


def _paginate_entries(entries, max_units=24):
    pages = []
    current = []
    current_units = 0

    for entry in entries:
        text = str(entry.get("text", "") or "").strip()
        if not text:
            continue
        chars_per_line = 30 if int(entry.get("font_size", 13) or 13) >= 13 else 36
        units = _estimate_line_units(text, chars_per_line=chars_per_line)
        units += 1 if int(entry.get("font_size", 13) or 13) >= 14 else 0
        if current and current_units + units > max_units:
            pages.append(current)
            current = []
            current_units = 0
        current.append(entry)
        current_units += units

    if current:
        pages.append(current)
    return pages or [[_make_entry("暂无详情", font_size=13)]]


def _fit_entries_to_single_page(entries, max_units=28, chars_per_line=34):
    adjusted = [dict(entry) for entry in entries if str(entry.get("text", "") or "").strip()]
    if not adjusted:
        return [_make_entry("暂无详情", font_size=13)]

    def total_units():
        dynamic_total = 0
        for entry in adjusted:
            font_size = int(entry.get("font_size", 13) or 13)
            dynamic_chars = chars_per_line + max(0, 13 - font_size) * 3
            dynamic_total += _entry_units(entry, chars_per_line=dynamic_chars)
        return dynamic_total

    if total_units() <= max_units:
        return adjusted

    removable_prefixes = ["原因:", "抓取概况:", "注意:", "⏱️ 承接时间线:"]
    while total_units() > max_units:
        removable_idx = next(
            (idx for idx in range(len(adjusted) - 1, -1, -1)
             if any(adjusted[idx]["text"].startswith(prefix) for prefix in removable_prefixes)),
            None,
        )
        if removable_idx is not None:
            adjusted.pop(removable_idx)
            continue
        break

    return adjusted


def _render_entries(text_frame, entries):
    text_frame.word_wrap = True
    for idx, entry in enumerate(entries):
        p = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        p.text = entry["text"]
        p.font.size = Pt(entry.get("font_size", 13))
        p.font.bold = bool(entry.get("bold", False))
        p.space_after = Pt(entry.get("space_after", 6))
        color = entry.get("color")
        if color:
            p.font.color.rgb = RGBColor(*color)
        hyperlink = entry.get("hyperlink")
        if hyperlink and p.runs:
            p.runs[0].hyperlink.address = hyperlink


def _build_news_entries(news, extraction_stats, section_warnings, compact=False, section_topic="", report_style=""):
    is_consumer_daily_section = report_style == "consumer_daily" and section_topic in _CONSUMER_DAILY_TOPIC_TITLES
    if is_consumer_daily_section:
        extraction_stats = {}
        section_warnings = []
    meta_font = 12
    stat_font = 9
    warn_font = 10
    body_font = 13
    title_ref_limit = 22 if compact else 28
    reason_limit = 58 if compact else 95
    timeline_ref_limit = 1 if compact else 2
    if is_consumer_daily_section:
        confidence = _get(news, "confidence_level", "") or "verified"
        source_count = _get(news, "independent_source_count", "") or "?"
        evidence_sources = _get(news, "evidence_sources", []) or []
        main_sources = "、".join(evidence_sources[:4]) if evidence_sources else _get(news, "source", "未知来源")
        confidence_note = "（待核实）" if confidence == "likely" else ""
        meta_text = (
            f"专题: {section_topic}  |  地区: {_guess_news_region(news)}  |  "
            f"时间: {_get(news, 'date_check', '近期')} / {_format_time_window_label(_get(news, 'event_time_window', ''))}  |  可信度: {confidence}{confidence_note}  |  "
            f"独立来源: {source_count} 个  |  主要来源: {_shorten(_clean_display_text(main_sources), 34)}"
        )
    else:
        meta_text = (
            (f"专题: {section_topic}  |  " if section_topic else "")
            + f"来源: {_get(news, 'source', '未知来源')}  |  时间: {_get(news, 'date_check', '近期')}  |  热度: {'⭐' * int(_get(news, 'importance', 3) or 3)}"
        )
    entries = [
        _make_entry(
            meta_text,
            font_size=meta_font,
            color=(128, 128, 128),
            space_after=6,
        )
    ]

    event_id = _get(news, "event_id", "")
    if event_id and not is_consumer_daily_section:
        entries.append(_make_entry(f"🧷 事件ID: {event_id}", font_size=stat_font + 1, color=(100, 100, 100), space_after=5))

    if extraction_stats:
        entries.append(
            _make_entry(
                f"抓取概况: {_format_extraction_stats(extraction_stats)}",
                font_size=stat_font,
                color=(100, 100, 100),
                space_after=4,
            )
        )

    if section_warnings:
        entries.append(
            _make_entry(
                f"注意: {section_warnings[0]}",
                font_size=warn_font,
                bold=True,
                color=(192, 102, 0),
                space_after=4,
            )
        )

    timeline_refs = _get(news, "timeline_refs", [])
    for ref in timeline_refs[:timeline_ref_limit]:
        entries.append(
            _make_entry(
                f"⏱️ 承接时间线: [{_get(ref, 'date', '近期')}] {_shorten(_clean_display_text(_get(ref, 'event', '')), title_ref_limit)}",
                font_size=stat_font + 1,
                bold=True,
                color=(192, 102, 0),
                space_after=3,
            )
        )
        if report_style != "company_tracking":
            entries.append(
                _make_entry(
                    f"原因: {_shorten(_clean_display_text(_get(ref, 'reason', '')), reason_limit)}",
                    font_size=stat_font,
                    color=(124, 45, 18),
                    space_after=4,
                )
            )

    _append_summary_entries(
        entries,
        _get(news, "summary", "暂无详情"),
        body_font=body_font,
        space_after=6,
    )

    if is_consumer_daily_section:
        evidence_urls = _get(news, "evidence_urls", []) or []
        if evidence_urls:
            entries.append(
                _make_entry(
                    f"证据链: {_shorten(' / '.join(evidence_urls[:3]), 120)}",
                    font_size=10,
                    color=(90, 90, 90),
                    space_after=4,
                )
            )

    news_url = _get(news, "url", "")
    if news_url:
        entries.append(
            _make_entry(
                "🔗 溯源查证: 点击查看原文",
                font_size=11,
                color=(0, 112, 192),
                hyperlink=news_url,
            )
        )

    return entries


def _add_chart_to_slide(slide, news, chart_info):
    chart_path = None
    try:
        chart_path = cg.generate_and_download_chart(
            _get(chart_info, "chart_title", "") or _get(news, "title", "数据图表"),
            _get(chart_info, "labels", []),
            _get(chart_info, "values", []),
            _get(chart_info, "chart_type", "bar"),
        )
    except Exception:
        chart_path = None

    if chart_path and os.path.exists(chart_path):
        slide.shapes.add_picture(chart_path, Inches(5.0), Inches(1.7), width=Inches(4.2))

    note_box = slide.shapes.add_textbox(Inches(5.0), Inches(6.25), Inches(4.2), Inches(0.55))
    note_p = note_box.text_frame.paragraphs[0]
    note_p.text = f"图表说明: {_shorten(_clean_display_text(_get(chart_info, 'chart_title', '') or _get(news, 'title', '数据对比')), 58)}"
    note_p.font.size = Pt(10)
    note_p.font.color.rgb = RGBColor(0, 51, 102)



def _format_extraction_stats(stats):
    stats = stats or {}
    return (
        f"Jina全文 {int(stats.get('jina_count', 0) or 0)} | "
        f"网页直连 {int(stats.get('direct_html_count', 0) or 0)} | "
        f"摘要兜底 {int(stats.get('snippet_count', 0) or 0)}"
    )


_CONSUMER_DAILY_TOPIC_TITLES = {
    "消费电子与手机新品",
    "消费电子 / 手机产业",
    "AR/VR与AI眼镜",
    "AR / VR / XR / AI 眼镜",
    "AI国内外重要资讯",
    "AI 一周资讯",
    "电动汽车智能科技",
    "电动汽车 / 智能汽车科技资讯",
    "折叠屏与新型显示",
    "折叠手机 / Fast LCD / LCoS / 显示与近眼显示供应链",
    "机器人 / 具身智能",
}


def _is_consumer_daily_deck(data):
    sections = list(data or [])
    if not sections:
        return False
    style_hits = sum(1 for section in sections if _get(section, "report_style", "") == "consumer_daily")
    topic_hits = sum(1 for section in sections if _get(section, "topic", "") in _CONSUMER_DAILY_TOPIC_TITLES)
    return style_hits >= 1 and topic_hits >= min(3, len(sections))


def _guess_news_region(news):
    blob = f"{_get(news, 'title', '')} {_get(news, 'summary', '')} {_get(news, 'source', '')}"
    if re.search(r"中国|国内|国产|华为|小米|荣耀|vivo|OPPO|比亚迪|小鹏|理想|蔚来|雷鸟|Rokid|XREAL|豆包|DeepSeek|通义|文心|混元", blob, re.IGNORECASE):
        return "中国/国内"
    return "海外/全球"


def _summary_one_sentence(news, limit=86):
    summary = _clean_display_text(_get(news, "summary", ""))
    summary = re.sub(r"【[^】]+】", "", summary).strip()
    if not summary:
        return "公开材料未披露更多摘要。"
    first = re.split(r"[。！？!?]\s*", summary)[0].strip()
    return _shorten(first or summary, limit)


def _format_time_window_label(value):
    key = str(value or "").lower()
    return {
        "today": "今日",
        "24h": "近24小时",
        "72h": "近72小时",
        "7d": "近一周",
    }.get(key, key or "当前窗口")


def _event_sources_text(event, limit=4):
    names = _get(event, "source_names", []) or []
    if names:
        return "、".join(str(item) for item in names[:limit])
    domains = _get(event, "source_domains", []) or []
    return "、".join(str(item) for item in domains[:limit]) if domains else "待补充来源"


def _add_consumer_daily_overview_slide(prs, sections):
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
    clear_placeholders(slide)

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(8.7), Inches(0.65))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "今日核心看点"
    title_p.font.size = Pt(26)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)

    total_news = sum(len(_get(section, "data", []) or []) for section in sections)
    meta_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(8.7), Inches(0.45))
    meta_p = meta_box.text_frame.paragraphs[0]
    meta_p.text = f"日期: {datetime.date.today()}  |  专题 {len(sections)} 个  |  详细新闻 {total_news} 条  |  中国国内主线优先，海外重大动态补充"
    meta_p.font.size = Pt(12)
    meta_p.font.color.rgb = RGBColor(100, 100, 100)

    body_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.8), Inches(8.25), Inches(4.8))
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, section in enumerate(sections):
        news_items = _get(section, "data", []) or []
        top_title = _shorten(_clean_display_text(_get(news_items[0], "title", "暂无重点新闻")) if news_items else "暂无重点新闻", 42)
        tags = "、".join((_get(section, "focus_tags", []) or [])[:4])
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"{idx + 1}. {_get(section, 'topic', '未命名专题')}：{len(news_items)} 条 | {tags} | {top_title}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = RGBColor(32, 32, 32)
        p.space_after = Pt(10)


def _extract_quality_report(sections):
    for section in sections or []:
        report = _get(section, "quality_report", {}) or {}
        if report:
            return report
    return {}


def _add_consumer_daily_quality_slide(prs, sections):
    report = _extract_quality_report(sections)
    if not report:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
    clear_placeholders(slide)

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(8.7), Inches(0.65))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "质量门槛与证据链"
    title_p.font.size = Pt(25)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)

    metrics = [
        f"正式事件: {_get(report, 'total_events', 0)} 条",
        f"confirmed: {_get(report, 'confirmed_events', 0)} 条",
        f"likely: {_get(report, 'likely_events', 0)} 条",
        f"weak/rejected: {_get(report, 'weak_events', 0)} / {_get(report, 'rejected_events', 0)} 条",
        f"来源多样性: {_get(report, 'source_diversity_score', 0)}",
    ]
    metric_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(8.4), Inches(1.1))
    metric_tf = metric_box.text_frame
    metric_tf.word_wrap = True
    for idx, item in enumerate(metrics):
        p = metric_tf.paragraphs[0] if idx == 0 else metric_tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.bold = idx < 3
        p.font.color.rgb = RGBColor(32, 32, 32)

    warnings = _get(report, "warnings", []) or []
    insufficient = _get(report, "topics_with_insufficient_events", []) or []
    stale = _get(report, "stale_events", []) or []
    topic_counts = _get(report, "topic_event_counts", {}) or {}
    confirmed_counts = _get(report, "topic_confirmed_counts", {}) or {}
    likely_counts = _get(report, "topic_likely_counts", {}) or {}
    watchlist_counts = _get(report, "topic_watchlist_counts", {}) or {}
    expansion_attempts = _get(report, "expansion_attempts", {}) or {}
    pipeline_stats = _get(report, "topic_pipeline_stats", {}) or {}
    lines = []
    for topic_name, count in list(topic_counts.items())[:5]:
        pipe = pipeline_stats.get(topic_name, {}) or {}
        pipeline_text = ""
        if pipe:
            pipeline_text = (
                f"，query {pipe.get('query_count', 0)}，URL {pipe.get('found_urls', 0)}，"
                f"时效保留 {pipe.get('freshness_kept', 0)}，抓取 {pipe.get('crawler_valid_count', 0)}，"
                f"事件主档 {pipe.get('event_master_count', 0)}"
            )
        lines.append(
            f"{topic_name}: 主新闻 {count} 条，confirmed {confirmed_counts.get(topic_name, 0)}，"
            f"likely {likely_counts.get(topic_name, 0)}，待跟踪 {watchlist_counts.get(topic_name, 0)}，"
            f"扩搜 {len(expansion_attempts.get(topic_name, []) or [])} 轮{pipeline_text}"
        )
    if warnings:
        lines.extend([f"警告: {item}" for item in warnings[:4]])
    if insufficient:
        lines.append("可确认新闻不足专题: " + "、".join(insufficient[:5]))
    if stale:
        lines.append("已剔除旧闻示例: " + "、".join(stale[:4]))
    if not lines:
        lines.append("本次正式正文仅使用 confirmed / likely 事件，已过滤单源、旧闻、主题错配和低质量聚合页。")

    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.75), Inches(8.4), Inches(3.8))
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(120, 70, 0) if "警告" in line or "不足" in line else RGBColor(32, 32, 32)
        p.space_after = Pt(10)


def _add_consumer_daily_list_slide(prs, section, section_index):
    news_items = _get(section, "data", []) or []
    if not news_items:
        return

    chunk_size = 3
    for offset in range(0, len(news_items), chunk_size):
        chunk = news_items[offset:offset + chunk_size]
        slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
        clear_placeholders(slide)

        title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9.0), Inches(0.65))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = f"{section_index}. {_get(section, 'topic', '未命名专题')} - 新闻列表"
        title_p.font.size = Pt(23)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(0, 51, 102)

        tags = "、".join((_get(section, "focus_tags", []) or [])[:7])
        if tags:
            tag_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(9.0), Inches(0.35))
            tag_p = tag_box.text_frame.paragraphs[0]
            tag_p.text = f"标签: {tags}"
            tag_p.font.size = Pt(10)
            tag_p.font.color.rgb = RGBColor(100, 100, 100)

        y = 1.48
        for local_idx, news in enumerate(chunk, start=offset + 1):
            title = _shorten(_clean_display_text(_get(news, "title", "未命名情报")), 46)
            confidence = _get(news, "confidence_level", "") or "verified"
            source_count = _get(news, "independent_source_count", "") or "?"
            evidence_sources = _get(news, "evidence_sources", []) or []
            main_sources = "、".join(evidence_sources[:3]) if evidence_sources else _get(news, "source", "未知来源")
            meta = (
                f"公司/产品: {title}  |  地区: {_guess_news_region(news)}  |  "
                f"时间: {_get(news, 'date_check', '近期')} / {_format_time_window_label(_get(news, 'event_time_window', ''))}  |  可信度: {confidence}  |  "
                f"独立来源: {source_count}  |  主要来源: {_shorten(_clean_display_text(main_sources), 24)}"
            )
            summary = _summary_one_sentence(news, limit=92)
            url = _get(news, "url", "")

            box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(8.6), Inches(1.45))
            tf = box.text_frame
            tf.word_wrap = True
            p_title = tf.paragraphs[0]
            p_title.text = f"{local_idx}. {title}"
            p_title.font.size = Pt(14)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(192, 102, 0)
            p_title.space_after = Pt(2)

            p_meta = tf.add_paragraph()
            p_meta.text = meta
            p_meta.font.size = Pt(9)
            p_meta.font.color.rgb = RGBColor(100, 100, 100)
            p_meta.space_after = Pt(2)

            p_summary = tf.add_paragraph()
            p_summary.text = f"摘要: {summary}"
            p_summary.font.size = Pt(10)
            p_summary.font.color.rgb = RGBColor(32, 32, 32)
            p_summary.space_after = Pt(2)

            if url:
                p_url = tf.add_paragraph()
                p_url.text = "来源 URL: 点击查看原文"
                p_url.font.size = Pt(9)
                p_url.font.color.rgb = RGBColor(0, 112, 192)
                if p_url.runs:
                    p_url.runs[0].hyperlink.address = url

            y += 1.75


def _add_consumer_daily_watchlist_slide(prs, section, section_index):
    watchlist = _get(section, "watchlist_events", []) or []
    if not watchlist:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
    clear_placeholders(slide)

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9.0), Inches(0.65))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = f"{section_index}. {_get(section, 'topic', '未命名专题')} - 待跟踪线索"
    title_p.font.size = Pt(23)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(120, 70, 0)

    reason = _get(section, "insufficient_reason", "") or "以下线索尚未满足多源验证标准，不能视为已确认新闻。"
    note_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.08), Inches(9.0), Inches(0.52))
    note_p = note_box.text_frame.paragraphs[0]
    note_p.text = f"说明: {reason}"
    note_p.font.size = Pt(11)
    note_p.font.bold = True
    note_p.font.color.rgb = RGBColor(120, 70, 0)

    body_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.75), Inches(8.55), Inches(4.9))
    tf = body_box.text_frame
    tf.word_wrap = True
    for idx, event in enumerate(watchlist[:2]):
        title = _shorten(_clean_display_text(_get(event, "normalized_title", "待跟踪线索")), 52)
        confidence = _get(event, "confidence_level", "weak")
        source_count = _get(event, "independent_source_count", 1)
        time_window = _format_time_window_label(_get(event, "time_window", ""))
        p_title = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p_title.text = f"{idx + 1}. {title}"
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(120, 70, 0)
        p_title.space_after = Pt(3)

        p_meta = tf.add_paragraph()
        p_meta.text = (
            f"状态: {confidence} / 待进一步核实  |  独立来源: {source_count}  |  "
            f"时间窗口: {time_window}  |  来源: {_shorten(_event_sources_text(event), 34)}"
        )
        p_meta.font.size = Pt(10)
        p_meta.font.color.rgb = RGBColor(100, 100, 100)
        p_meta.space_after = Pt(3)

        p_summary = tf.add_paragraph()
        p_summary.text = "线索摘要: " + _shorten(_clean_display_text(_get(event, "event_summary", "公开材料未披露更多摘要。")), 125)
        p_summary.font.size = Pt(11)
        p_summary.font.color.rgb = RGBColor(32, 32, 32)
        p_summary.space_after = Pt(12)


def _add_consumer_daily_closing_slide(prs, sections):
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
    clear_placeholders(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.65), Inches(8.6), Inches(0.75))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = "后续跟踪清单"
    title_p.font.size = Pt(26)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)

    watch_terms = []
    for section in sections:
        watch_terms.extend((_get(section, "watch_entities", []) or [])[:5])
    deduped_terms = []
    seen = set()
    for item in watch_terms:
        value = str(item or "").strip()
        if not value or value in seen:
            continue
        seen.add(value)
        deduped_terms.append(value)
        if len(deduped_terms) >= 18:
            break

    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(8.2), Inches(4.8))
    tf = body_box.text_frame
    tf.word_wrap = True
    lines = [
        "趋势判断: 继续跟踪硬件参数、量产节奏、供应链订单、价格和渠道变化，避免只看发布会口径。",
        f"重点公司/产品: {'、'.join(deduped_terms) if deduped_terms else '公开材料未披露'}",
        "待验证信息: 爆料类参数、未披露销量、未确认供应商、融资与量产时间表需要后续来源交叉验证。",
        "下一轮搜索重点: 中国国内科技媒体、公司官方新闻稿、供应链垂直媒体和海外重大产品发布。",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16 if idx == 0 else 14)
        p.font.bold = idx == 0
        p.font.color.rgb = RGBColor(32, 32, 32)
        p.space_after = Pt(14)


def clear_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)


def _add_section_divider_slide(prs, section, section_index, section_total):
    slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
    clear_placeholders(slide)

    kicker_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.15), Inches(8.5), Inches(0.35))
    kicker_p = kicker_box.text_frame.paragraphs[0]
    kicker_p.text = f"专题模块 {section_index}/{section_total}"
    kicker_p.font.size = Pt(14)
    kicker_p.font.color.rgb = RGBColor(100, 100, 100)
    kicker_p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.85), Inches(8.5), Inches(1.1))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = _get(section, "topic", "未命名专题")
    title_p.font.size = Pt(30)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(0, 51, 102)
    title_p.alignment = PP_ALIGN.CENTER

    focus_tags = _get(section, "focus_tags", [])
    watch_entities = _get(section, "watch_entities", [])
    meta_parts = []
    if focus_tags:
        meta_parts.append(f"重点标签: {'、'.join(focus_tags[:8])}")
    if watch_entities:
        meta_parts.append(f"关注主体: {'、'.join(watch_entities[:8])}")

    if meta_parts:
        meta_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.25), Inches(8.0), Inches(1.2))
        meta_tf = meta_box.text_frame
        meta_tf.word_wrap = True
        for idx, line in enumerate(meta_parts):
            p = meta_tf.paragraphs[0] if idx == 0 else meta_tf.add_paragraph()
            p.text = line
            p.font.size = Pt(15)
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.alignment = PP_ALIGN.CENTER

    count_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.35), Inches(8.0), Inches(0.6))
    count_p = count_box.text_frame.paragraphs[0]
    count_p.text = f"详细新闻 {len(_get(section, 'data', []) or [])} 条"
    count_p.font.size = Pt(16)
    count_p.font.bold = True
    count_p.font.color.rgb = RGBColor(192, 102, 0)
    count_p.alignment = PP_ALIGN.CENTER



def generate_ppt(data, timeline_data, filename, model_name):
    template_path = "template.pptx"
    if os.path.exists(template_path):
        try:
            prs = Presentation(template_path)
        except Exception:
            prs = Presentation()
    else:
        prs = Presentation()

    all_sections = list(data or [])
    sections_with_news = [section for section in all_sections if _get(section, "data", [])]
    is_consumer_daily = _is_consumer_daily_deck(all_sections)
    sections_for_output = [
        section for section in all_sections
        if _get(section, "data", []) or (is_consumer_daily and _get(section, "watchlist_events", []))
    ]

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_box.text_frame.paragraphs[0].text = "《科技消费电子日报》" if is_consumer_daily else "《FPC-RD 科技资讯》"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_box.text_frame.paragraphs[0].text = (
        f"生成日期: {datetime.date.today()}  |  中国国内科技新闻优先"
        if is_consumer_daily else
        f"生成日期: {datetime.date.today()}"
    )
    subtitle_box.text_frame.paragraphs[0].font.size = Pt(16)
    subtitle_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if is_consumer_daily:
        if sections_with_news:
            _add_consumer_daily_overview_slide(prs, sections_with_news)
        _add_consumer_daily_quality_slide(prs, all_sections)

    if timeline_data and not is_consumer_daily:
        for t_data in timeline_data:
            events = _get(t_data, "events", [])
            if not events:
                continue

            timeline_report_style = _get(t_data, "report_style", "")
            is_company_tracking_timeline = timeline_report_style == "company_tracking"
            chunk_size = 3 if is_company_tracking_timeline else 5
            for i in range(0, len(events), chunk_size):
                chunk = events[i:i + chunk_size]
                slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
                clear_placeholders(slide)

                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(9), Inches(0.8))
                title_box.text_frame.paragraphs[0].text = f"⏱️ {_get(t_data, 'topic', '未命名专题')} - 核心时间线"
                title_box.text_frame.paragraphs[0].font.size = Pt(24)
                title_box.text_frame.paragraphs[0].font.bold = True

                meta_parts = []
                focus_tags = _get(t_data, "focus_tags", [])
                if focus_tags:
                    meta_parts.append(f"重点标签: {'、'.join(focus_tags[:6])}")
                extraction_stats = _get(t_data, "extraction_stats", {})
                if extraction_stats:
                    meta_parts.append(_format_extraction_stats(extraction_stats))
                if meta_parts:
                    meta_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.35))
                    meta_p = meta_box.text_frame.paragraphs[0]
                    meta_p.text = " | ".join(meta_parts)
                    meta_p.font.size = Pt(10)
                    meta_p.font.color.rgb = RGBColor(100, 100, 100)

                body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(9), Inches(5.45))
                tf = body_box.text_frame
                tf.word_wrap = True
                for idx, item in enumerate(chunk):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    appears_later = bool(_get(item, "appears_in_later_news", False))
                    history_status = _get(item, "history_status", "")
                    prefix = "★ " if appears_later else ""
                    if history_status == "followup":
                        prefix = "◆ " + prefix
                    display_event = _shorten(_clean_display_text(_get(item, "event", "未命名事件")), 42)
                    display_source = _shorten(_clean_display_text(_get(item, "source", "未知来源")), 18)
                    p.text = f"{prefix}[{_get(item, 'date', '近期')}] {display_event} ({display_source})"
                    p.font.size = Pt(14)
                    p.space_after = Pt(2)
                    if appears_later:
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(192, 102, 0)
                    elif history_status == "followup":
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(15, 118, 110)

                    if is_company_tracking_timeline:
                        summary_text = _shorten(_format_timeline_event_summary(item), 230)
                        if summary_text:
                            p_summary = tf.add_paragraph()
                            p_summary.text = summary_text
                            p_summary.font.size = Pt(11)
                            p_summary.font.color.rgb = RGBColor(31, 78, 121)
                            p_summary.space_after = Pt(4)

                            source_url = str(_get(item, "source_url", "") or "").strip()
                            if source_url:
                                p_link = tf.add_paragraph()
                                p_link.space_after = Pt(4)
                                label_run = p_link.add_run()
                                label_run.text = "🔗 原新闻："
                                label_run.font.size = Pt(9)
                                label_run.font.color.rgb = RGBColor(0, 112, 192)
                                link_run = p_link.add_run()
                                link_run.text = "点击查看原文"
                                link_run.font.size = Pt(9)
                                link_run.font.color.rgb = RGBColor(0, 112, 192)
                                link_run.hyperlink.address = source_url

                    if history_status == "followup":
                        p_hist = tf.add_paragraph()
                        p_hist.text = (
                            f"  ↳ 历史追踪: 首次记录 {_get(item, 'first_seen', '未知')} / "
                            f"累计 {int(_get(item, 'seen_count', 1) or 1)} 次"
                        )
                        p_hist.font.size = Pt(10)
                        p_hist.font.color.rgb = RGBColor(15, 118, 110)
                        p_hist.space_after = Pt(5)

                    if appears_later:
                        p_reason = tf.add_paragraph()
                        if is_company_tracking_timeline:
                            matched_title = _shorten(_clean_display_text(_get(item, "matched_news_title", "")), 42)
                            if not matched_title:
                                matched_title = "相关详细新闻"
                            p_reason.text = f"  ↳ 详见后文：《{matched_title}》"
                        else:
                            p_reason.text = (
                                f"  ↳ 在《{_shorten(_clean_display_text(_get(item, 'matched_news_title', '')), 32)}》中展开："
                                f"{_shorten(_clean_display_text(_get(item, 'match_reason', '')), 90)}"
                            )
                        p_reason.font.size = Pt(10)
                        p_reason.font.color.rgb = RGBColor(124, 45, 18)
                        p_reason.space_after = Pt(8)

    for section_index, section in enumerate(sections_for_output, start=1):
        news_items = _get(section, "data", [])
        if not news_items and not (is_consumer_daily and _get(section, "watchlist_events", [])):
            continue

        if len(sections_for_output) > 1:
            _add_section_divider_slide(prs, section, section_index, len(sections_for_output))
        if is_consumer_daily:
            _add_consumer_daily_list_slide(prs, section, section_index)
            _add_consumer_daily_watchlist_slide(prs, section, section_index)

        finance = _get(section, "finance", {}) or {}
        section_warnings = _get(section, "warnings", [])
        extraction_stats = _get(section, "extraction_stats", {})
        focus_tags = _get(section, "focus_tags", [])
        report_style = _get(section, "report_style", "")

        if finance.get("is_public"):
            finance_slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
            clear_placeholders(finance_slide)

            title_box = finance_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            title_box.text_frame.paragraphs[0].text = f"📊 {_get(section, 'topic', '未命名专题')} ({finance.get('ticker', '')}) - 量化面与事件催化"
            title_box.text_frame.paragraphs[0].font.size = Pt(22)
            title_box.text_frame.paragraphs[0].font.bold = True

            if focus_tags or extraction_stats:
                meta_box = finance_slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.3))
                meta_p = meta_box.text_frame.paragraphs[0]
                meta_parts = []
                if focus_tags:
                    meta_parts.append(f"重点标签: {'、'.join(focus_tags[:6])}")
                if extraction_stats:
                    meta_parts.append(_format_extraction_stats(extraction_stats))
                meta_p.text = " | ".join(meta_parts)
                meta_p.font.size = Pt(9)
                meta_p.font.color.rgb = RGBColor(100, 100, 100)

            body_box = finance_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.0), Inches(2.8))
            tf = body_box.text_frame

            data_available = finance.get("data_available", True)
            if not data_available:
                msg = finance.get("msg", "Financial data temporarily unavailable.")
                p_msg = tf.paragraphs[0]
                p_msg.text = msg
                p_msg.font.size = Pt(14)
                p_msg.font.color.rgb = RGBColor(128, 128, 128)
            else:
                change_pct_raw = finance.get("change_pct")
                change_pct = _to_float(change_pct_raw)
                if change_pct is None:
                    trend_icon = "⏺"
                    color = RGBColor(90, 90, 90)
                    change_pct_text = "N/A"
                elif change_pct > 0:
                    trend_icon = "🔺"
                    color = RGBColor(192, 0, 0)
                    change_pct_text = f"{change_pct:.2f}"
                elif change_pct < 0:
                    trend_icon = "🔻"
                    color = RGBColor(0, 150, 0)
                    change_pct_text = f"{change_pct:.2f}"
                else:
                    trend_icon = "⏺"
                    color = RGBColor(90, 90, 90)
                    change_pct_text = "0.00"

                current_price = _fmt_number(finance.get("current_price"))
                currency = finance.get("currency", "")
                p_price = tf.paragraphs[0]
                p_price.text = f"{current_price} {currency}  {trend_icon} {change_pct_text}%"
                p_price.font.size = Pt(24)
                p_price.font.bold = True
                p_price.font.color.rgb = color

                metrics = [
                    f"▪ 估值水平: {finance.get('pe_pb', 'N/A')}",
                    f"▪ 股权风险溢价: {finance.get('erp', 'N/A')}",
                    f"▪ 总市值: {finance.get('market_cap', 'N/A')}",
                ]
                for metric in metrics:
                    p = tf.add_paragraph()
                    p.text = metric
                    p.font.size = Pt(13)
                    p.space_before = Pt(12)

                chart_path = finance.get("chart_path")
                if chart_path and os.path.exists(chart_path):
                    finance_slide.shapes.add_picture(chart_path, Inches(4.5), Inches(1.2), width=Inches(5.0))

            catalysts = finance.get("catalysts", {})
            boxes_data = [
                ("🏛️ 政策与监管", catalysts.get("policy", "近期无重大政策催化")),
                ("💰 财报与盈利", catalysts.get("earnings", "未见核心财报数据")),
                ("🚀 产业标志事件", catalysts.get("landmark", "产业层级平稳")),
                ("🔄 市场风格轮动", catalysts.get("style", "风格未见明显切换")),
            ]

            for i, (title, content) in enumerate(boxes_data):
                x_pos = 0.5 + (i * 2.2)
                content_box = finance_slide.shapes.add_textbox(Inches(x_pos), Inches(4.5), Inches(2.1), Inches(2.5))
                content_tf = content_box.text_frame
                content_tf.word_wrap = True
                p_title = content_tf.paragraphs[0]
                p_title.text = title
                p_title.font.size = Pt(12)
                p_title.font.bold = True
                p_title.font.color.rgb = RGBColor(0, 51, 102)
                p_content = content_tf.add_paragraph()
                p_content.text = content
                p_content.font.size = Pt(11)
                p_content.space_before = Pt(6)

            if extraction_stats:
                stat_box = finance_slide.shapes.add_textbox(Inches(0.5), Inches(6.65), Inches(9.0), Inches(0.3))
                stat_p = stat_box.text_frame.paragraphs[0]
                stat_p.text = f"抓取概况: {_format_extraction_stats(extraction_stats)}"
                stat_p.font.size = Pt(9)
                stat_p.font.color.rgb = RGBColor(100, 100, 100)
            if section_warnings:
                warn_box = finance_slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9.0), Inches(0.5))
                warn_p = warn_box.text_frame.paragraphs[0]
                warn_p.text = f"注意: {section_warnings[0]}"
                warn_p.font.size = Pt(10)
                warn_p.font.bold = True
                warn_p.font.color.rgb = RGBColor(192, 102, 0)

        for news in news_items:
            chart_info = _get(news, "chart_info", {})
            has_chart = bool(_get(chart_info, "has_chart", False) and len(_get(chart_info, "labels", [])) > 0)
            entries = _build_news_entries(
                news,
                extraction_stats,
                section_warnings,
                compact=has_chart,
                section_topic=_get(section, "topic", ""),
                report_style=report_style,
            )

            if has_chart:
                fitted_entries = _fit_entries_to_single_page(entries, max_units=24, chars_per_line=23)
                if _total_units(fitted_entries, chars_per_line=23) > 24:
                    has_chart = False
                    entries = _build_news_entries(
                        news,
                        extraction_stats,
                        section_warnings,
                        compact=False,
                        section_topic=_get(section, "topic", ""),
                        report_style=report_style,
                    )
            if not has_chart:
                fitted_entries = _fit_entries_to_single_page(entries, max_units=29, chars_per_line=34)

            slide = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
            clear_placeholders(slide)

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
            title_text = _shorten(_get(news, "title", "未命名情报"), 38 if has_chart else 56)
            title_box.text_frame.paragraphs[0].text = title_text
            title_box.text_frame.paragraphs[0].font.size = Pt(22)
            title_box.text_frame.paragraphs[0].font.bold = True

            if has_chart:
                body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.45), Inches(4.2), Inches(5.4))
                _render_entries(body_box.text_frame, fitted_entries)
                _add_chart_to_slide(slide, news, chart_info)
            else:
                body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9.0), Inches(5.8))
                _render_entries(body_box.text_frame, fitted_entries)

    if is_consumer_daily and sections_with_news:
        _add_consumer_daily_closing_slide(prs, sections_with_news)

    path = f"{filename}.pptx"
    prs.save(path)
    return path
