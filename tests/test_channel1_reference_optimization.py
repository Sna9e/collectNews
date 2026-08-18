import pathlib
import re
import sys
import tempfile

from pptx import Presentation
from pptx.dml.color import RGBColor


ROOT = pathlib.Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from agents.deep_analyst import NewsItem, _backfill_event_ids  # noqa: E402
from tools.company_query_packs import (  # noqa: E402
    DEFAULT_COMPANY_TOPICS,
    build_company_focus_hint,
    build_company_queries_from_pack,
    get_company_query_pack,
)
from tools.export_ppt import generate_ppt  # noqa: E402
from tools.report_linker import annotate_report_data  # noqa: E402


def _event_summary(subject, action, object_text, detail, impact):
    return (
        f"{subject}于8月18日{action}，本次动作直接围绕{object_text}展开。"
        f"相关方案覆盖{detail}，并明确了产品或业务的实施范围。"
        f"现有材料还披露了客户验证、交付节奏或供应链准备等具体节点。"
        f"这将直接影响{impact}，相关团队需要据此调整后续执行计划。"
    )


def _build_link_fixture():
    deep_sections = [
        {
            "topic": "OpenAI",
            "report_style": "company_tracking",
            "data": [
                {
                    "event_id": "E99",
                    "title": "OpenAI发布新一代企业模型与API",
                    "source": "OpenAI",
                    "date_check": "08月18日",
                    "url": "https://openai.com/index/model-update?utm_source=test",
                    "importance": 5,
                    "summary": (
                        "【事件核心】OpenAI发布新一代企业模型与API，面向开发者和企业客户开放。"
                        "【深度细节/数据支撑】更新涉及模型能力、接口调用和企业部署。"
                        "【行业深远影响】该动作将影响企业智能体产品的开发节奏。"
                    ),
                    "chart_info": {"has_chart": False},
                },
                {
                    "event_id": "E02",
                    "title": "苹果调整新一代iPhone供应链备货",
                    "source": "Apple Newsroom",
                    "date_check": "08月18日",
                    "url": "https://apple.com/newsroom/iphone-supply",
                    "importance": 3,
                    "summary": (
                        "【事件核心】苹果调整新一代iPhone供应链备货安排。"
                        "【深度细节/数据支撑】变化涉及显示模组、主板和摄像头零部件。"
                        "【行业深远影响】上游厂商需要重新安排排产、库存和交付窗口。"
                    ),
                    "chart_info": {"has_chart": False},
                },
            ],
            "finance": {},
            "warnings": [],
            "extraction_stats": {},
        }
    ]
    timeline_sections = [
        {
            "topic": "OpenAI",
            "report_style": "company_tracking",
            "events": [
                {
                    "event_id": "E01",
                    "date": "08月18日",
                    "source": "OpenAI",
                    "source_url": "https://www.openai.com/index/model-update/",
                    "event": "OpenAI发布企业模型API",
                    "event_summary": _event_summary(
                        "OpenAI", "发布新一代企业模型与API", "模型能力和接口调用",
                        "开发者开放、企业部署和智能体工具", "企业客户部署与开发者适配节奏",
                    ),
                },
                {
                    "event_id": "E02",
                    "date": "08月18日",
                    "source": "Apple Newsroom",
                    "source_url": "",
                    "event": "苹果调整iPhone供应链备货",
                    "event_summary": _event_summary(
                        "苹果", "调整新一代iPhone供应链备货", "显示模组、主板和摄像头零部件",
                        "排产窗口、库存安排和交付准备", "上游供应商的生产与库存计划",
                    ),
                },
                {
                    "event_id": "E03",
                    "date": "08月18日",
                    "source": "其他媒体",
                    "source_url": "",
                    "event": "苹果零部件订单节奏变化",
                    "event_summary": _event_summary(
                        "苹果", "调整部分零部件订单节奏", "新一代iPhone供应链",
                        "显示、主板和摄像头环节", "供应商排产与交付安排",
                    ),
                },
            ],
        }
    ]
    return annotate_report_data(deep_sections, timeline_sections)


def test_all_ten_topics_have_specific_content_profiles():
    assert DEFAULT_COMPANY_TOPICS == [
        "Apple", "Google", "Amazon", "OpenAI", "Meta",
        "Nvidia", "Tesla", "特朗普", "Anthropic", "SpaceX",
    ]
    for topic in DEFAULT_COMPANY_TOPICS:
        pack = get_company_query_pack(topic)
        queries = build_company_queries_from_pack(topic, pack)
        hint = build_company_focus_hint(pack)
        assert pack["id"] != "generic"
        assert len(queries) >= 8
        assert len(pack.get("summary_focus", [])) >= 4
        assert len(pack.get("high_value_signals", [])) >= 6
        assert any(re.search(r"[\u4e00-\u9fff]", query) for query in queries)
        assert any(re.search(r"[A-Za-z]", query) for query in queries)
        assert "短新闻摘要和详细新闻应优先交代" in hint
        assert "用于重点高亮" in hint


def test_deep_news_event_id_backfill_prefers_normalized_source_url():
    item = NewsItem(
        event_id="E02",
        title="OpenAI发布企业模型API",
        source="OpenAI",
        date_check="08月18日",
        url="https://openai.com/index/model-update?utm_campaign=daily",
        summary="【事件核心】OpenAI发布企业模型API。",
    )
    blueprints = [
        {
            "event_id": "E01",
            "event": "OpenAI发布企业模型API",
            "source_url": "https://www.openai.com/index/model-update/",
            "keywords": ["OpenAI", "API"],
        },
        {
            "event_id": "E02",
            "event": "OpenAI调整办公室安排",
            "source_url": "https://example.com/office",
            "keywords": ["办公室"],
        },
    ]
    _backfill_event_ids([item], blueprints)
    assert item.event_id == "E01"


def test_report_linker_uses_url_then_validated_event_id_and_is_one_to_one():
    deep_sections, timeline_sections = _build_link_fixture()
    news = deep_sections[0]["data"]
    events = timeline_sections[0]["events"]

    assert events[0]["match_method"] == "source_url"
    assert events[0]["matched_news_index"] == 1
    assert events[0]["highlight_level"] == "key"
    assert events[1]["match_method"] == "event_id"
    assert events[1]["matched_news_index"] == 2
    assert events[1]["highlight_level"] == "linked"
    assert events[2]["appears_in_later_news"] is False
    assert [item["detail_index"] for item in news] == [1, 2]
    assert [len(item["timeline_refs"]) for item in news] == [1, 1]


def test_ppt_shows_numbered_mapping_and_two_highlight_levels():
    deep_sections, timeline_sections = _build_link_fixture()
    with tempfile.TemporaryDirectory() as tmpdir:
        ppt_path = generate_ppt(
            deep_sections,
            timeline_sections,
            str(pathlib.Path(tmpdir) / "channel1_reference_optimization"),
            "stub",
        )
        presentation = Presentation(ppt_path)

        all_text = []
        timeline_paragraphs = []
        for slide in presentation.slides:
            slide_text = "\n".join(
                shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
            )
            all_text.append(slide_text)
            if "OpenAI - 核心时间线" in slide_text:
                for shape in slide.shapes:
                    if not getattr(shape, "has_text_frame", False):
                        continue
                    timeline_paragraphs.extend(shape.text_frame.paragraphs)

        joined = "\n".join(all_text)
        assert "详见后文：详细新闻 1《OpenAI发布新一代企业模型与API》" in joined
        assert "详见后文：详细新闻 2《苹果调整新一代iPhone供应链备货》" in joined
        assert "详细新闻 1｜OpenAI发布新一代企业模型与API" in joined
        assert "详细新闻 2｜苹果调整新一代iPhone供应链备货" in joined
        assert "事件ID" not in joined
        assert "原因:" not in joined

        key_title = next(p for p in timeline_paragraphs if p.text.startswith("★ [08月18日]"))
        linked_title = next(p for p in timeline_paragraphs if p.text.startswith("● [08月18日]"))
        assert key_title.font.color.rgb == RGBColor(192, 102, 0)
        assert linked_title.font.color.rgb == RGBColor(31, 78, 121)


def run_all():
    tests = [
        test_all_ten_topics_have_specific_content_profiles,
        test_deep_news_event_id_backfill_prefers_normalized_source_url,
        test_report_linker_uses_url_then_validated_event_id_and_is_one_to_one,
        test_ppt_shows_numbered_mapping_and_two_highlight_levels,
    ]
    for test in tests:
        test()
        print(f"PASS {test.__name__}")


if __name__ == "__main__":
    run_all()
