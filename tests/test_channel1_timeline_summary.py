import pathlib
import re
import sys
import tempfile

from pptx import Presentation
from pptx.util import Pt


ROOT = pathlib.Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from agents.timeline_agent import (  # noqa: E402
    _build_event_summary_from_result,
    _merge_event_dict,
    build_event_blueprints,
    generate_timeline,
)
from tools.export_ppt import generate_ppt  # noqa: E402


APPLE_STUB_PPT = ROOT / "stub_validation_channel1_timeline_apple.pptx"


def _cjk_count(text):
    return len(re.findall(r"[\u4e00-\u9fff]", str(text or "")))


def _sentence_count(text):
    return len(re.findall(r"[。！？]", str(text or "")))


def _estimate_lines(text, chars_per_line=42):
    weighted = 0.0
    for ch in str(text or ""):
        if "\u4e00" <= ch <= "\u9fff":
            weighted += 1.0
        elif ord(ch) < 128:
            weighted += 0.55
        else:
            weighted += 0.8
    return max(1, int((weighted + chars_per_line - 1) // chars_per_line))


def test_normal_search_content_builds_complete_summary():
    result = {
        "content": (
            "英伟达宣布 Vera Rubin 平台进入全面生产，面向下一代 AI 工厂提供 CPU、GPU 和网络组件。"
            "报道称该平台将用于提升数据中心训练与推理效率，并承接 Blackwell 之后的产品节奏。"
            "材料提到服务器厂商和云服务客户会围绕新平台推进验证。"
            "供应链环节因此需要同步准备板卡、散热和网络配套。"
            "短期影响是数据中心客户的采购规划更集中于新平台切换。"
        )
    }
    summary = _build_event_summary_from_result(result, {"event": "英伟达 Vera Rubin 投产"})
    assert summary
    assert "英伟达宣布 Vera Rubin 平台进入全面生产" in summary
    assert 140 <= len(summary) <= 220
    assert 4 <= _sentence_count(summary) <= 5


def test_summary_cleaner_removes_noise_and_patch_lines():
    result = {
        "content": """
网页导航 分享 订阅 newsletter privacy policy
进一步看，第二条摘要不应进入时间线。
补充判断：围绕英伟达还可以继续观察。
英伟达宣布 Vera Rubin 平台进入全面生产，相关材料指向下一代 AI 工厂建设和数据中心产品节奏。
报道提到该平台覆盖 CPU、GPU 和网络组件，服务数据中心训练与推理任务。
供应链需要围绕服务器板卡、散热和网络配套同步准备。
该动作会影响云服务客户的新平台验证安排，也会改变下一阶段采购节奏。
相关阅读 热门推荐 read more
"""
    }
    summary = _build_event_summary_from_result(result, {"event": "英伟达 Vera Rubin 投产"})
    assert "英伟达宣布 Vera Rubin 平台进入全面生产" in summary
    assert "进一步看" not in summary
    assert "补充判断：围绕" not in summary
    assert "newsletter" not in summary.lower()
    assert "read more" not in summary.lower()


def test_empty_material_returns_no_summary():
    summary = _build_event_summary_from_result({}, {"event": ""})
    assert summary == ""


def test_english_result_does_not_replace_chinese_event_summary():
    result = {
        "content": "NVIDIA confirms Vera Rubin launch for Q3 as Blackwell demand keeps climbing [...] Read more",
        "title": "NVIDIA confirms Vera Rubin launch for Q3",
    }
    summary = _build_event_summary_from_result(result, {"event": "英伟达确认Vera Rubin芯片", "keywords": ["Vera Rubin", "Blackwell"]})
    assert summary == ""
    assert "NVIDIA confirms" not in summary
    assert "Read more" not in summary
    assert "[...]" not in summary


def test_disclaimer_summary_is_rejected():
    result = {
        "content": (
            "公开材料显示，苹果发布端侧AI计划。该线索于06月01日由某网站披露。"
            "材料没有提供足够细节，暂不能确认更多参数。时间线仅记录已披露动作。"
        )
    }
    summary = _build_event_summary_from_result(result, {"event": "Apple发布端侧AI计划"})
    assert summary == ""


def test_merge_keeps_better_summary_without_concatenating():
    existing = {
        "event": "英伟达芯片进展",
        "event_summary": "",
        "keywords": ["GPU"],
    }
    candidate = {
        "event": "英伟达芯片进展",
        "event_summary": "英伟达披露新平台生产进展，材料明确提到数据中心和 AI 工厂相关产品节奏。报道提到服务器厂商会围绕新平台推进验证。云服务客户的训练与推理部署计划也会随之调整。该动作会影响板卡、散热和网络配套的供应链准备节奏。",
        "keywords": ["AI"],
    }
    merged = _merge_event_dict(existing, candidate)
    assert merged["event_summary"] == candidate["event_summary"]
    assert merged["event_summary"].count("英伟达") == 1


def _ppt_text_by_slide(path):
    slides = []
    for slide in Presentation(path).slides:
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text)
        slides.append("\n".join(parts))
    return slides


def _timeline_paragraphs(path, topic="Apple"):
    rows = []
    for slide_index, slide in enumerate(Presentation(path).slides, start=1):
        slide_text = "\n".join(
            shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
        )
        if f"{topic} - 核心时间线" not in slide_text:
            continue
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    rows.append({"slide": slide_index, "text": text, "font_size": paragraph.font.size})
    return rows


def _apple_events():
    return [
        {
            "date": "06月01日",
            "event": "Apple发布端侧AI计划",
            "source": "Apple Newsroom",
            "source_url": "https://www.apple.com/newsroom/apple-ai",
            "event_summary": "苹果在6月发布端侧AI计划，重点面向iPhone和Mac提供本地模型能力。该计划覆盖系统体验、开发者接口和设备端推理场景。报道提到芯片算力、隐私处理和应用适配会成为落地重点。该动作将影响后续系统升级节奏，也会提高终端硬件对AI能力的配置要求。供应链和开发者需要据此调整测试与适配安排。",
            "appears_in_later_news": True,
            "matched_news_title": "苹果端侧AI计划推动设备体验升级",
            "match_reason": "这是一段不应在 PPT 核心时间线中展示的很长匹配原因。",
        },
        {
            "date": "06月01日",
            "event": "iPhone供应链备货调整",
            "source": "供应链日报",
            "source_url": "https://example.com/apple-supply-chain",
            "event_summary": "供应链消息称苹果调整新一代iPhone备货节奏，涉及摄像头、显示模组和主板环节。公开材料指向部分零部件订单节奏变化，但未披露具体数量。该变化会影响上游厂商排产、库存和交付安排。对苹果而言，备货调整有助于匹配新品发布前后的需求预期，并减少发布初期供需错配风险。供应链企业需要根据后续订单确认生产窗口。",
        },
        {
            "date": "06月02日",
            "event": "Mac芯片路线更新",
            "source": "Developer News",
            "source_url": "https://example.com/apple-mac-chip",
            "event_summary": "苹果更新Mac芯片路线，公开信息显示新平台聚焦能效、图形性能和AI推理能力。材料提到高端笔记本和桌面设备会优先受益。该路线会改变开发者对本地推理、图形渲染和电池续航的优化重点，并影响专业软件生态的性能验证计划。对硬件业务而言，新芯片有助于强化Mac产品差异化。后续节点取决于具体机型发布安排。",
            "history_status": "followup",
            "first_seen": "05月20日",
            "seen_count": 2,
        },
        {
            "date": "06月02日",
            "event": "Vision Pro渠道调整",
            "source": "Retail Watch",
            "source_url": "https://example.com/apple-vision-pro",
            "event_summary": "苹果在6月调整Vision Pro渠道安排，相关动作涉及线下展示、零售培训和重点市场覆盖。报道没有披露具体门店数量，但指向体验式销售资源重新分配。渠道调整会影响消费者试用路径，也会改变开发者和内容合作方观察设备需求的节奏。后续销售表现将影响苹果对空间计算硬件的资源投入。",
        },
    ]


def test_ppt_timeline_summaries_pagination_and_link_prompt():
    events = _apple_events()
    data = [
        {
            "topic": "Apple",
            "data": [
                {
                    "title": "苹果端侧AI计划推动设备体验升级",
                    "source": "Apple Newsroom",
                    "date_check": "06月01日",
                    "summary": "【事件核心】\n测试。\n【深度细节/数据支撑】\n测试。\n【行业深远影响】\n测试。",
                    "importance": 3,
                    "chart_info": {"has_chart": False},
                }
            ],
            "report_style": "company_tracking",
            "finance": {},
            "warnings": [],
            "extraction_stats": {},
            "focus_tags": [],
        }
    ]
    timeline_data = [{"topic": "Apple", "events": events, "report_style": "company_tracking"}]

    ppt_path = generate_ppt(data, timeline_data, str(APPLE_STUB_PPT.with_suffix("")), "stub")
    slide_texts = _ppt_text_by_slide(ppt_path)

    timeline_slides = [text for text in slide_texts if "核心时间线" in text]
    assert len(timeline_slides) == 2
    for text in timeline_slides:
        assert len(re.findall(r"(?m)^\s*(?:★\s*|◆\s*)?\[\d{2}月\d{2}日\]", text)) <= 3

    all_timeline_text = "\n".join(timeline_slides)
    assert "端侧AI计划" in all_timeline_text
    assert "供应链消息称苹果调整新一代iPhone备货节奏" in all_timeline_text
    assert "苹果更新Mac芯片路线" in all_timeline_text
    assert "↳ 详见后文：《苹果端侧AI计划推动设备体验升级》" in all_timeline_text
    assert "很长匹配原因" not in all_timeline_text
    assert "中展开" not in all_timeline_text
    assert "Apple launches" not in all_timeline_text
    assert "进一步看" not in all_timeline_text
    assert "补充判断" not in all_timeline_text
    assert "公开材料显示" not in all_timeline_text
    assert "该线索" not in all_timeline_text
    assert "材料没有提供足够" not in all_timeline_text

    paragraphs = _timeline_paragraphs(ppt_path, topic="Apple")
    by_text = {row["text"]: row for row in paragraphs}
    normal_summaries = [event["event_summary"] for event in events if event["event_summary"]]
    for summary in normal_summaries:
        assert summary in by_text
        assert 100 <= len(summary) <= 220
        assert 3 <= _sentence_count(summary) <= 5
        assert _cjk_count(summary) >= 80
        assert by_text[summary]["font_size"] in {Pt(10), Pt(11)}
    for slide_text in timeline_slides:
        assert sum(_estimate_lines(line) for line in slide_text.splitlines() if line.strip()) <= 28


class FakeAI:
    valid = True

    def analyze_structural(self, prompt, structure_class):
        if structure_class.__name__ == "EventBlueprintReport":
            assert "event_summary" in prompt
            assert "3 到 5 个完整句子" in prompt
            assert "100 到 220" in prompt
            assert "公开材料显示" in prompt
            return structure_class(
                events=[
                    {
                        "date": "06月01日",
                        "source": "Apple Newsroom",
                        "event": "Apple发布端侧AI计划",
                        "event_summary": "苹果在6月发布端侧AI计划，面向iPhone和Mac提供本地模型能力。该计划会覆盖系统体验、开发者接口和设备端推理场景。相关内容提到芯片算力、隐私处理和应用适配是落地重点。该动作会影响系统升级节奏，也会提高终端硬件对AI能力的配置要求。开发者需要围绕新接口调整测试与适配安排。",
                        "source_url": "https://www.apple.com/newsroom/apple-ai",
                        "keywords": ["Apple", "端侧AI", "iPhone"],
                    }
                ]
            )
        if structure_class.__name__ == "TimelineTitleReport":
            return structure_class(events=[{"event": "Apple发布端侧AI计划"}])
        raise AssertionError(structure_class.__name__)


def test_fake_ai_blueprint_summary_flows_to_timeline_and_ppt():
    raw_results = [
        {
            "title": "Apple launches on-device AI plan for iPhone and Mac",
            "content": "Apple launches on-device AI plan for iPhone and Mac. Read more",
            "url": "https://www.apple.com/newsroom/apple-ai",
            "source": "Apple Newsroom",
            "published_at_resolved": "2026-06-01T10:00:00+00:00",
        }
    ]
    blueprints = build_event_blueprints(FakeAI(), raw_results, "Apple", "2026年06月05日", "过去1个月")
    timeline = generate_timeline(blueprints)
    assert timeline[0].event_summary.startswith("苹果在6月发布端侧AI计划")
    assert "Apple launches" not in timeline[0].event_summary

    data = [{"topic": "Apple", "data": [], "report_style": "company_tracking", "finance": {}}]
    timeline_data = [{"topic": "Apple", "events": timeline, "report_style": "company_tracking"}]
    with tempfile.TemporaryDirectory() as tmpdir:
        ppt_path = generate_ppt(data, timeline_data, str(pathlib.Path(tmpdir) / "fake_ai_flow"), "stub")
        all_text = "\n".join(_ppt_text_by_slide(ppt_path))
    assert "苹果在6月发布端侧AI计划，面向iPhone和Mac提供本地模型能力" in all_text


def test_company_tracking_marker_is_only_in_channel1_block():
    text = (ROOT / "agent_app.py").read_text(encoding="utf-8")
    start = text.index("with tab1:")
    tab2 = text.index("with tab2:", start)
    tab3 = text.index("with tab3:", tab2)
    channel1_block = text[start:tab2]
    channel2_block = text[tab2:tab3]
    assert '"report_style": "company_tracking"' in channel1_block
    assert '"report_style": "company_tracking"' not in channel2_block


def run_all():
    tests = [
        test_normal_search_content_builds_complete_summary,
        test_summary_cleaner_removes_noise_and_patch_lines,
        test_empty_material_returns_no_summary,
        test_english_result_does_not_replace_chinese_event_summary,
        test_disclaimer_summary_is_rejected,
        test_merge_keeps_better_summary_without_concatenating,
        test_ppt_timeline_summaries_pagination_and_link_prompt,
        test_fake_ai_blueprint_summary_flows_to_timeline_and_ppt,
        test_company_tracking_marker_is_only_in_channel1_block,
    ]
    for test in tests:
        test()
        print(f"PASS {test.__name__}")


if __name__ == "__main__":
    run_all()
