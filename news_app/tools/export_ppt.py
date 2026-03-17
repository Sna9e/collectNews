# -*- coding: utf-8 -*-
import datetime
import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from news_app.tools import chart_generator as cg


def clear_placeholders(slide):
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            sp = shape.element
            sp.getparent().remove(sp)


def _safe_float(val, default=0.0):
    try:
        return float(val)
    except Exception:
        return default


def _template_path() -> str:
    root = Path(__file__).resolve().parents[2]
    return str(root / "assets" / "template.pptx")


def generate_ppt(data, timeline_data, filename, model_name):
    template_path = _template_path()
    if os.path.exists(template_path):
        try:
            prs = Presentation(template_path)
        except Exception:
            prs = Presentation()
    else:
        prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    clear_placeholders(slide)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_box.text_frame.paragraphs[0].text = "高管情报：前沿动态深度分析"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_box.text_frame.paragraphs[0].text = (
        f"生成日期: {datetime.date.today()} | 模型: {model_name}"
    )
    subtitle_box.text_frame.paragraphs[0].font.size = Pt(16)
    subtitle_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    if timeline_data:
        for t_data in timeline_data:
            if not t_data["events"]:
                continue
            chunk_size = 7
            events = t_data["events"]
            for i in range(0, len(events), chunk_size):
                chunk = events[i : i + chunk_size]
                slide = prs.slides.add_slide(
                    prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                )
                clear_placeholders(slide)

                t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
                t_box.text_frame.paragraphs[0].text = f"时间线：{t_data['topic']}"
                t_box.text_frame.paragraphs[0].font.size = Pt(24)
                t_box.text_frame.paragraphs[0].font.bold = True

                b_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
                tf = b_box.text_frame
                tf.word_wrap = True
                for idx, item in enumerate(chunk):
                    p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                    p.text = f"[{item.date}] {item.event} ({item.source})"
                    p.font.size = Pt(14)
                    p.space_after = Pt(8)

    for section in data:
        if not section["data"]:
            continue

        finance = section.get("finance", {})
        if finance.get("is_public") and finance.get("data_ok"):
            f_slide = prs.slides.add_slide(
                prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            )
            clear_placeholders(f_slide)

            t_box = f_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
            t_box.text_frame.paragraphs[0].text = (
                f"量化概览：{section['topic']} ({finance.get('ticker', '')})"
            )
            t_box.text_frame.paragraphs[0].font.size = Pt(22)
            t_box.text_frame.paragraphs[0].font.bold = True

            b_box = f_slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.0), Inches(2.8))
            tf = b_box.text_frame

            change_pct = _safe_float(finance.get("change_pct"))
            trend_icon = "↑" if change_pct > 0 else "↓"
            color = RGBColor(192, 0, 0) if change_pct > 0 else RGBColor(0, 150, 0)
            p_price = tf.paragraphs[0]
            p_price.text = (
                f"{finance.get('current_price')} {finance.get('currency')}  "
                f"{trend_icon} {change_pct:.2f}%"
            )
            p_price.font.size = Pt(24)
            p_price.font.bold = True
            p_price.font.color.rgb = color

            metrics = [
                f"估值: {finance.get('pe_pb', 'N/A')}",
                f"股权风险溢价: {finance.get('erp', 'N/A')}",
                f"市值: {finance.get('market_cap', 'N/A')}",
            ]
            for m in metrics:
                p = tf.add_paragraph()
                p.text = m
                p.font.size = Pt(13)
                p.space_before = Pt(12)

            chart_path = finance.get("chart_path")
            if chart_path and os.path.exists(chart_path):
                f_slide.shapes.add_picture(chart_path, Inches(4.5), Inches(1.2), width=Inches(5.0))

            cat = finance.get("catalysts", {})
            boxes_data = [
                ("政策与监管", cat.get("policy", "暂无明确催化")),
                ("财报与盈利", cat.get("earnings", "暂无关键财务信号")),
                ("产业事件", cat.get("landmark", "暂无标志性事件")),
                ("风格轮动", cat.get("style", "暂无明显切换")),
            ]

            for i, (title, content) in enumerate(boxes_data):
                x_pos = 0.5 + (i * 2.2)
                c_box = f_slide.shapes.add_textbox(Inches(x_pos), Inches(4.5), Inches(2.1), Inches(2.5))
                c_tf = c_box.text_frame
                c_tf.word_wrap = True

                p_t = c_tf.paragraphs[0]
                p_t.text = title
                p_t.font.size = Pt(12)
                p_t.font.bold = True
                p_t.font.color.rgb = RGBColor(0, 51, 102)

                p_c = c_tf.add_paragraph()
                p_c.text = content
                p_c.font.size = Pt(11)
                p_c.space_before = Pt(6)

        for news in section["data"]:
            slide = prs.slides.add_slide(
                prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            )
            clear_placeholders(slide)

            t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
            t_box.text_frame.paragraphs[0].text = news.title
            t_box.text_frame.paragraphs[0].font.size = Pt(22)
            t_box.text_frame.paragraphs[0].font.bold = True

            has_chart = (
                hasattr(news, "chart_info")
                and news.chart_info.has_chart
                and len(news.chart_info.labels) > 0
            )
            text_width = 5.2 if has_chart else 9.0

            b_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(text_width), Inches(5))
            tf = b_box.text_frame
            tf.word_wrap = True

            stars = "*" * max(1, int(news.importance))
            tf.paragraphs[0].text = (
                f"来源: {news.source}  |  时间: {news.date_check}  |  热度: {stars}"
            )
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
            tf.paragraphs[0].space_after = Pt(10)

            for line in news.summary.split("\n"):
                line = line.strip()
                if not line:
                    continue
                p = tf.add_paragraph()
                p.text = line
                p.font.size = Pt(13)
                p.space_after = Pt(6)
                if line.startswith("【") or line.startswith("["):
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0, 51, 102)

            news_url = getattr(news, "url", "")
            if news_url:
                p_link = tf.add_paragraph()
                p_link.text = "来源核验：点击查看原文"
                p_link.font.size = Pt(11)
                p_link.font.color.rgb = RGBColor(0, 112, 192)
                p_link.runs[0].hyperlink.address = news_url

            if has_chart:
                try:
                    chart_img = cg.generate_and_download_chart(
                        news.chart_info.chart_title,
                        news.chart_info.labels,
                        news.chart_info.values,
                        news.chart_info.chart_type,
                    )
                    if chart_img and os.path.exists(chart_img):
                        slide.shapes.add_picture(
                            chart_img, Inches(5.8), Inches(1.5), width=Inches(3.8)
                        )
                        try:
                            os.remove(chart_img)
                        except Exception:
                            pass
                except Exception:
                    pass

    path = f"{filename}.pptx"
    prs.save(path)
    return path
