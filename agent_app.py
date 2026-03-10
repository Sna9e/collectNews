import streamlit as st
import asyncio
import sys
import json
import os
import datetime
import subprocess
import concurrent.futures
import platform
import urllib.request
import urllib.parse
import difflib
from typing import List

# ================= 0. 核心库引用 =================
from pydantic import BaseModel, Field, ValidationError
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openai import OpenAI  

# ================= 1. 核心网络配置 =================
if platform.system() == "Windows":
    os.environ["http_proxy"] = "http://127.0.0.1:7890"
    os.environ["https_proxy"] = "http://127.0.0.1:7890"
else:
    os.environ.pop("http_proxy", None)
    os.environ.pop("https_proxy", None)
    os.environ.pop("HTTP_PROXY", None)
    os.environ.pop("HTTPS_PROXY", None)

# ================= 2. 爬虫、文档与PPT排版库引用 =================
from crawl4ai import AsyncWebCrawler
from docx import Document
from docx.shared import RGBColor, Pt as DocxPt
from docx.oxml.ns import qn
from docx.enum.text import WD_ALIGN_PARAGRAPH 

# 🔴 新增：PPT 核心排版引擎
from pptx import Presentation
from pptx.util import Inches, Pt as PptxPt

if sys.platform == 'win32':
    asyncio.set_event_loop_policy(asyncio.WindowsProactorEventLoopPolicy())

st.set_page_config(page_title="DeepSeek 科技探员", page_icon="🐳", layout="wide")

# ================= 3. 定义结构化数据 =================
class NewsItem(BaseModel):
    title: str = Field(description="新闻标题（务必翻译为中文）")
    source: str = Field(description="来源媒体（保留原名）")
    date_check: str = Field(description="严格核实新闻发生的真实日期，格式 YYYY-MM-DD。")
    summary: str = Field(description="约300字的深度商业分析。必须严格分段并带有标识：【事件核心】、【深度细节/数据支撑】、【行业深远影响】。")
    importance: int = Field(description="重要性 1-5")

class NewsReport(BaseModel):
    news: List[NewsItem] = Field(description="新闻列表")

# ================= 4. 内置 DeepSeek 驱动 =================
class EnterpriseDeepSeekDriver:
    def __init__(self, api_key, model_id):
        self.valid = False
        if not api_key: return
        try:
            self.client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com")
            self.model_id = model_id
            self.valid = True
        except Exception:
            pass

    def analyze_structural(self, prompt, structure_class):
        if not self.valid: return None
        schema_str = json.dumps(structure_class.model_json_schema(), ensure_ascii=False)
        sys_prompt = f"你是顶级商业情报分析师。必须严格按此 JSON Schema 返回数据，不带任何废话：\n{schema_str}"
        try:
            response = self.client.chat.completions.create(
                model=self.model_id,
                messages=[{"role": "system", "content": sys_prompt}, {"role": "user", "content": prompt}],
                response_format={"type": "json_object"},
                temperature=0.1, 
                max_tokens=4096 
            )
            raw_text = response.choices[0].message.content.strip()
            try:
                json_obj = json.loads(raw_text)
                if isinstance(json_obj, list): json_obj = {"news": json_obj}
                return structure_class(**json_obj)
            except Exception: return None
        except Exception: return None

# ================= 5. 核心业务函数 =================

@st.cache_resource(show_spinner="☁️ 首次启动：正在云端配置无头浏览器内核 (约需1-2分钟)...")
def check_and_install_playwright():
    try:
        subprocess.run([sys.executable, "-m", "playwright", "install", "chromium"], check=True, capture_output=True)
        if platform.system() != "Windows":
            subprocess.run([sys.executable, "-m", "playwright", "install-deps", "chromium"], check=True, capture_output=True)
        return True
    except Exception:
        return False

def search_web(query, sites_text, timelimit, max_results=10, tavily_key=""):
    if not tavily_key: return []
    sites = [s.strip() for s in sites_text.split('\n') if s.strip()]
    
    try:
        url = "https://api.tavily.com/search"
        payload = {
            "api_key": tavily_key,
            "query": query, 
            "search_depth": "advanced",
            "topic": "news", 
            "max_results": max_results
        }
        if sites: payload["include_domains"] = sites
        
        if timelimit == "d": payload["days"] = 2 
        elif timelimit == "w": payload["days"] = 7
        elif timelimit == "m": payload["days"] = 30

        data = json.dumps(payload).encode('utf-8')
        req = urllib.request.Request(url, data=data, headers={'Content-Type': 'application/json'})
        resp = json.loads(urllib.request.urlopen(req, timeout=15).read().decode('utf-8'))
        
        results = [{'href': r['url']} for r in resp.get('results', [])]
        return results
    except Exception as e:
        print(f"Tavily Search Failed: {e}")
        return []

async def crawl_urls_concurrently(urls):
    full_content = ""
    valid_count = 0
    async with AsyncWebCrawler() as crawler:
        tasks = [crawler.arun(url=url) for url in urls]
        results = await asyncio.gather(*tasks, return_exceptions=True)
        for i, res in enumerate(results):
            if isinstance(res, Exception): continue
            if res.success:
                valid_count += 1
                markdown_text = res.fit_markdown if hasattr(res, 'fit_markdown') and res.fit_markdown else res.markdown
                if markdown_text and len(markdown_text) > 200:
                    full_content += f"\n\n=== SOURCE START: {urls[i]} ===\n{markdown_text[:6000]}\n=== SOURCE END ===\n"
    return full_content, valid_count

def safe_run_async_crawler(urls):
    new_loop = asyncio.new_event_loop()
    asyncio.set_event_loop(new_loop)
    try:
        return new_loop.run_until_complete(crawl_urls_concurrently(urls))
    finally:
        new_loop.close()

def map_reduce_analysis(ai_driver, topic, full_text, current_date, time_opt):
    if not full_text or len(full_text) < 100: return []
    docs = RecursiveCharacterTextSplitter(chunk_size=8000, chunk_overlap=1000).create_documents([full_text])
    all_extracted_news = []

    def process_single_doc(doc):
        map_prompt = f"""
        【全局时间锚点】：今天是 **{current_date}**。
        要求的时间范围是：【{time_opt}】。
        任务：从以下文本提取关于【{topic}】的新闻情报。
        红线：
        1. 严格时间审查：发现发生时间早于【{time_opt}】之前（如几个月前、或者去年），直接丢弃相关新闻，但是深度细节中不用提到丢弃了新闻！
        2. 【{topic}】必须是绝对主角！
        无符合条件的内容必须返回 `{{"news": []}}`。
        文本：{doc.page_content}
        """
        return ai_driver.analyze_structural(map_prompt, NewsReport)

    with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
        for future in concurrent.futures.as_completed([executor.submit(process_single_doc, d) for d in docs]):
            res = future.result()
            if res and res.news: all_extracted_news.extend(res.news)

    if not all_extracted_news: return []
    combined_json = json.dumps([item.model_dump() for item in all_extracted_news], ensure_ascii=False)

    reduce_prompt = f"""
        【全局时间锚点】：今天是 **{current_date}**。
        你是极其严苛的科技媒体总编。
        任务：
        1. 终极时间清洗：任何陈年旧闻，全部无情删掉！
        2. 合并去重：报道同一事件的新闻必须合并。
        3. 深度扩写与高级排版：将每条新闻的 summary 扩展至 300 字左右。必须在 summary 中使用明显的分段和换行，明确包含以下三个部分：
           【事件核心】：概括事件，将整体事件清晰明确的描述出来
           【深度细节】：核心数据与细节支撑，需要把核心细节与行业关注数据详细摘录。不用特意提到在新闻在时效性内
           【行业影响】：精简的行业深远影响，此事件可能对行业以及相关产业有什么影响，市场会对此有什么反应与措施
        4. 按重要性降序，最多保留最核心的 5 条。
        数据：{combined_json}
    """
    final_report = ai_driver.analyze_structural(reduce_prompt, NewsReport)
    return final_report.news if final_report else []

def generate_word(data, filename, model_name):
    doc = Document()
    normal_style = doc.styles['Normal']
    normal_style.font.name = '微软雅黑'
    normal_style._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
    normal_style.font.size = DocxPt(10.5) 
    
    for i in range(1, 4):
        h_style = doc.styles[f'Heading {i}']
        h_style.font.name = '微软雅黑'
        h_style._element.rPr.rFonts.set(qn('w:eastAsia'), '微软雅黑')
        if i == 1:
            h_style.font.color.rgb = RGBColor(0, 51, 102)

    title = doc.add_heading("DeepSeek 企业级深度科技研报", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    meta_p = doc.add_paragraph()
    meta_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta_run = meta_p.add_run(f"生成日期: {datetime.date.today()}  |  数据来源: Tavily 商业资讯引擎  |  分析模型: {model_name}")
    meta_run.font.color.rgb = RGBColor(128, 128, 128)
    meta_run.font.size = DocxPt(9)
    
    doc.add_paragraph("━" * 50).alignment = WD_ALIGN_PARAGRAPH.CENTER

    for section in data:
        doc.add_heading(f"🔷 专题：{section['topic']}", level=1)
        if not section['data']:
            doc.add_paragraph("    在指定时间范围内，未发现符合标准的重大情报。").font.italic = True
            continue
            
        for news in section['data']:
            doc.add_heading(f"🔹 {news.title}", level=2)
            
            p_info = doc.add_paragraph()
            run_info = p_info.add_run(f"    📌 来源: {news.source}    |    🕒 时间: {news.date_check}    |    🔥 价值评级: {'⭐'*news.importance}")
            run_info.font.color.rgb = RGBColor(100, 100, 100)
            run_info.font.bold = True
            
            p_summary = doc.add_paragraph(news.summary)
            p_summary.paragraph_format.line_spacing = 1.5 
            p_summary.paragraph_format.first_line_indent = DocxPt(21) 
            
            divider = doc.add_paragraph("┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈┈")
            divider.alignment = WD_ALIGN_PARAGRAPH.CENTER
            divider.runs[0].font.color.rgb = RGBColor(200, 200, 200)
    
    path = f"{filename}.docx"
    doc.save(path)
    return path

# 🔴 新增：极速原生 PPT 生成器
def generate_ppt(data, filename, model_name):
    prs = Presentation()
    
    # 1. 制作高逼格封面
    title_slide_layout = prs.slide_layouts[0] # 封面排版
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "行业前沿情报深度分析"
    subtitle.text = f"生成日期: {datetime.date.today()}\n数据引擎: Tavily & {model_name}"

    # 2. 循环写入专题和新闻
    for section in data:
        # 添加【专题过渡页】
        if section['data']:
            section_layout = prs.slide_layouts[2] # 节标题排版
            sec_slide = prs.slides.add_slide(section_layout)
            sec_slide.shapes.title.text = f"🎯 追踪目标：{section['topic']}"
            
            # 为该专题下的每一条新闻生成一页 PPT
            for news in section['data']:
                content_layout = prs.slide_layouts[1] # 标题+内容排版
                slide = prs.slides.add_slide(content_layout)
                
                # 设置单页标题（新闻标题）
                title_shape = slide.shapes.title
                title_shape.text = news.title
                title_shape.text_frame.paragraphs[0].font.size = PptxPt(28)
                
                # 设置单页正文
                body_shape = slide.shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.clear() # 清空默认格式
                
                # 写入元数据 (来源、时间、热度)
                p_meta = tf.add_paragraph()
                p_meta.text = f"📌 来源: {news.source}  |  🕒 {news.date_check}  |  🔥 热度: {'⭐'*news.importance}"
                p_meta.font.size = PptxPt(14)
                p_meta.font.color.rgb = RGBColor(128, 128, 128)
                
                # 写入深度分析正文
                p_summary = tf.add_paragraph()
                p_summary.text = f"\n{news.summary}"
                p_summary.font.size = PptxPt(16)
    
    path = f"{filename}.pptx"
    prs.save(path)
    return path

# ================= 6. 主界面 =================
with st.sidebar:
    st.header("🐳 DeepSeek 控制台")
    api_key = st.text_input("DeepSeek API Key", type="password")
    
    tavily_key = st.text_input("Tavily API Key (必填)", type="password", help="必须填入此 Key 才能驱动云端极速新闻引擎！")
    
    model_id = st.selectbox("模型", ["deepseek-chat"], index=0)
    st.divider()
    
    time_opt = st.selectbox("时间范围（绝对严控）", ["过去 24 小时", "过去 1 周", "过去 1 个月", "不限时间"], index=0)
    time_limit_dict = {"过去 24 小时": "d", "过去 1 周": "w", "过去 1 个月": "m", "不限时间": None}
    
    st.markdown("**搜外媒请用英文名 (如 Google、Apple)**")
    sites = st.text_area("重点搜索源", "techcrunch.com\ntheverge.com\nengadget.com\ncnet.com\nbloomberg.com/technology\nelectrek.co\ninsideevs.com\nroadtovr.com\nuploadvr.com\n36kr.com\nithome.com\nhuxiu.com\ngeekpark.net\nvrtuoluo.cn\nd1ev.com", height=250)
    file_name = st.text_input("文件名", f"深度研报_{datetime.date.today()}")

st.title("🐳 企业情报探员 (PPT双擎输出版)")
query_input = st.text_input("输入主题 (用 \\ 隔开，外媒源建议用英文如：Google \\ Apple)", "Google \\ OpenAI \\ Anthropic")
btn = st.button("🚀 开始生成研报", type="primary")

if btn:
    if not api_key or not tavily_key:
        st.error("❌ 请先在左侧边栏填入 DeepSeek 和 Tavily 的 API Key！")
    elif not query_input:
        st.warning("请输入关键词！")
    else:
        check_and_install_playwright()
        
        topics = [t.strip() for t in query_input.split('\\') if t.strip()]
        all_data = []
        ai = EnterpriseDeepSeekDriver(api_key, model_id)
        current_date_str = datetime.date.today().strftime("%Y年%m月%d日")
        
        global_seen_titles = []

        st.info("🚀 探员已出击，Tavily 企业级检索正在运行...")

        for topic in topics:
            st.markdown(f"#### 🔵 追踪目标: 【{topic}】 (要求: {time_opt})")
            
            with st.spinner(f"正在全网搜寻关于【{topic}】的最新线索..."):
                links = search_web(topic, sites, time_limit_dict[time_opt], tavily_key=tavily_key)
            
            if not links: 
                st.warning(f"⚠️ {topic}：在严格的【{time_opt}】限制内未搜寻到新闻。说明目标近期很安静！")
                continue
                
            st.write(f"🔍 成功截获 {len(links)} 个相关网址，启动智能爬虫...")

            with st.spinner(f"正在并发抓取并提纯这 {len(links)} 个网页的正文..."):
                full_text_data, valid_count = safe_run_async_crawler(urls=[r['href'] for r in links])

            if full_text_data:
                st.write(f"🧠 成功提取 {valid_count} 个纯净网页。DeepSeek 正在执行清洗与排版归纳...")
                
                with st.spinner("AI 正在剔除旧闻与重复项，撰写商业分析..."):
                    final_news_list = map_reduce_analysis(ai, topic, full_text_data, current_date_str, time_opt)
                
                if final_news_list:
                    deduped_news = []
                    for news in final_news_list:
                        is_duplicate = False
                        for seen_title in global_seen_titles:
                            similarity = difflib.SequenceMatcher(None, news.title, seen_title).ratio()
                            if similarity > 0.6:
                                is_duplicate = True
                                break
                        
                        if not is_duplicate:
                            deduped_news.append(news)
                            global_seen_titles.append(news.title)
                    
                    if deduped_news:
                        all_data.append({"topic": topic, "data": deduped_news})
                        filtered_count = len(final_news_list) - len(deduped_news)
                        st.success(f"✅ 【{topic}】分析完毕！已锁定 {len(deduped_news)} 条新鲜情报。" + 
                                   (f"(已跨主题去重过滤 {filtered_count} 条重复事件)" if filtered_count > 0 else ""))
                    else:
                        st.warning(f"⚠️ 【{topic}】提炼出的新闻均与之前的主题高度重合，已执行全局去重抹杀！")
                        
                else:
                    st.warning(f"⚠️ 【{topic}】搜到的网页经 AI 严格审判，全被判定为旧闻或非核心新闻，已执行抹杀过滤。")
            else:
                st.error(f"❌ 网页抓取失败或正文均为空，目标网站反爬拦截。")
            
            st.divider()

        if all_data:
            # 🔴 同时生成 Word 和 PPT
            path_word = generate_word(all_data, file_name, model_id)
            path_ppt = generate_ppt(all_data, file_name, model_id)
            
            st.balloons()
            st.success("🎉 全链条任务执行完毕！老板专供版 PPT 已就绪。")
            
            # 🔴 并排显示两个下载按钮
            col1, col2 = st.columns(2)
            with col1:
                with open(path_word, "rb") as f:
                    st.download_button("📝 立即下载深度研报 (Word)", f, file_name=path_word, type="secondary")
            with col2:
                with open(path_ppt, "rb") as f:
                    # PPT 按钮设为 primary，吸引眼球
                    st.download_button("📊 立即下载汇报演示 (PPT)", f, file_name=path_ppt, type="primary")
        else:
            st.error(f"❌ 任务结束。在严格的时效与去重约束下，所有关键词均未产生独立且有效的大事件情报。")
